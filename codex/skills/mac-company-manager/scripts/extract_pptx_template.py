#!/usr/bin/env python3
"""
extract_pptx_template.py — MaC Company Manager
Extracts template specs from a PPTX file and generates:
  - manifest.json       (complete field inventory per slide)
  - content-schema.json (JSON Schema for content input)
  - copy-prompt.md      (content generation guide)

Usage:
  python extract_pptx_template.py <path-to.pptx> [--output-dir DIR] [--template-name NAME]

Requires: python-pptx
  pip install python-pptx
"""

import argparse
import json
import math
import os
import shutil
import sys
from pathlib import Path


def check_dependency():
    try:
        from pptx import Presentation  # noqa: F401
    except ImportError:
        print("Error: python-pptx is required.")
        print("Install it with: pip install python-pptx")
        sys.exit(1)


def emu_to_inches(emu):
    """Convert EMU (English Metric Units) to inches. 914400 EMU = 1 inch."""
    return round(emu / 914400, 3)


def estimate_max_chars(width_in, height_in, font_size_pt):
    """
    Estimate max characters that fit in a text box.
    Formula: chars_per_line * lines
      chars_per_line ≈ (width_in * 72) / (font_size_pt * 0.55)
      lines ≈ (height_in * 72) / (font_size_pt * 1.2)
    The 0.55 factor approximates average character width relative to em-size.
    The 1.2 factor approximates line height (120% of font size).
    """
    if font_size_pt <= 0:
        font_size_pt = 12
    chars_per_line = (width_in * 72) / (font_size_pt * 0.55)
    lines = (height_in * 72) / (font_size_pt * 1.2)
    return max(10, int(chars_per_line * lines * 0.85))  # 0.85 safety margin


def classify_field_type(shape_name, placeholder_type, text_content, font_size_pt):
    """Classify a text field as title, subtitle, body, caption, or label."""
    name_lower = shape_name.lower()
    text_lower = (text_content or "").lower()

    if placeholder_type is not None:
        from pptx.enum.text import PP_ALIGN  # noqa: F401
        from pptx.util import Pt  # noqa: F401
        # Placeholder type integers: 1=center_title, 2=body, 3=center_title,
        # 13=subtitle, 15=title
        if placeholder_type in (1, 3, 15):
            return "title"
        if placeholder_type == 13:
            return "subtitle"
        if placeholder_type == 2:
            return "body" if font_size_pt <= 24 else "subtitle"

    if any(k in name_lower for k in ("title", "heading")):
        return "title"
    if any(k in name_lower for k in ("subtitle", "sub_title")):
        return "subtitle"
    if any(k in name_lower for k in ("caption", "footnote", "source", "credit")):
        return "caption"
    if any(k in name_lower for k in ("label", "eyebrow", "tag", "kicker")):
        return "label"
    if font_size_pt and font_size_pt >= 28:
        return "title" if font_size_pt >= 36 else "subtitle"
    return "body"


def get_font_size(shape):
    """Extract font size from a shape's first paragraph/run, in points."""
    try:
        for para in shape.text_frame.paragraphs:
            if para.runs:
                sz = para.runs[0].font.size
                if sz:
                    return round(sz / 12700, 1)  # EMU to pt
            if para.font and para.font.size:
                return round(para.font.size / 12700, 1)
    except (AttributeError, IndexError):
        pass
    return 12.0  # default


def get_placeholder_type(shape):
    """Return placeholder type index or None if not a placeholder."""
    try:
        return shape.placeholder_format.type
    except AttributeError:
        return None


def process_slide(slide, slide_index, slide_width_in, slide_height_in):
    """Extract all fields from a single slide."""
    fields = []
    field_id_counters = {}

    layout_name = "Unknown"
    try:
        layout_name = slide.slide_layout.name
    except AttributeError:
        pass

    for shape in slide.shapes:
        field = None

        # --- Text shapes ---
        if shape.has_text_frame:
            text_content = shape.text_frame.text.strip()
            font_size = get_font_size(shape)
            ph_type = get_placeholder_type(shape)

            left = emu_to_inches(shape.left or 0)
            top = emu_to_inches(shape.top or 0)
            width = emu_to_inches(shape.width or int(slide_width_in * 914400))
            height = emu_to_inches(shape.height or int(0.5 * 914400))

            max_chars = estimate_max_chars(width, height, font_size)
            field_type = classify_field_type(shape.name, ph_type, text_content, font_size)

            base_id = field_type
            count = field_id_counters.get(base_id, 0)
            field_id_counters[base_id] = count + 1
            field_id = base_id if count == 0 else f"{base_id}_{count}"

            field = {
                "id": field_id,
                "type": "text",
                "label": shape.name,
                "position": {
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height,
                },
                "font_size_pt": font_size,
                "max_chars": max_chars,
                "required": field_type in ("title",),
                "example": text_content[:200] if text_content else "",
            }

        # --- Image/picture shapes ---
        elif shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            left = emu_to_inches(shape.left or 0)
            top = emu_to_inches(shape.top or 0)
            width = emu_to_inches(shape.width or int(2 * 914400))
            height = emu_to_inches(shape.height or int(1.5 * 914400))

            rec_w = int(width * 300)
            rec_h = int(height * 300)

            count = field_id_counters.get("image", 0)
            field_id_counters["image"] = count + 1
            field_id = "image" if count == 0 else f"image_{count}"

            field = {
                "id": field_id,
                "type": "image",
                "label": shape.name,
                "position": {
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height,
                },
                "recommended_px": {"width": rec_w, "height": rec_h},
                "required": False,
                "example": "",
            }

        # --- Table shapes ---
        elif shape.has_table:
            table = shape.table
            rows = len(table.rows)
            cols = len(table.columns)
            has_header = rows > 1

            count = field_id_counters.get("table", 0)
            field_id_counters["table"] = count + 1
            field_id = "table" if count == 0 else f"table_{count}"

            # Extract table data as list of row lists
            table_data = []
            for row in table.rows:
                row_data = [cell.text_frame.text.strip() for cell in row.cells]
                table_data.append(row_data)

            field = {
                "id": field_id,
                "type": "table",
                "label": shape.name,
                "rows": rows,
                "columns": cols,
                "has_header_row": has_header,
                "required": False,
                "example": table_data,
            }

        if field:
            fields.append(field)

    return {
        "slide_index": slide_index,
        "layout_name": layout_name,
        "fields": fields,
    }


def build_manifest(pptx_path, template_name, slides_data, slide_width_in, slide_height_in):
    return {
        "template_name": template_name,
        "source_file": os.path.basename(pptx_path),
        "slide_dimensions": {
            "width_in": round(slide_width_in, 3),
            "height_in": round(slide_height_in, 3),
        },
        "total_slides": len(slides_data),
        "slides": slides_data,
    }


def build_content_schema(manifest):
    """Generate JSON Schema from manifest."""
    properties = {}
    required = []

    for slide in manifest["slides"]:
        idx = slide["slide_index"]
        for field in slide["fields"]:
            if field["type"] in ("text",):
                key = f"slide_{idx}_{field['id']}"
                prop = {
                    "type": "string",
                    "description": (
                        f"Slide {idx + 1} {field['label']} — "
                        f"{field['font_size_pt']}pt, "
                        f"fits ~{field['max_chars']} chars"
                    ),
                }
                if field.get("max_chars"):
                    prop["maxLength"] = field["max_chars"]
                if field.get("example"):
                    prop["examples"] = [field["example"]]
                properties[key] = prop
                if field.get("required"):
                    required.append(key)

            elif field["type"] == "image":
                key = f"slide_{idx}_{field['id']}"
                rec = field.get("recommended_px", {})
                properties[key] = {
                    "type": "string",
                    "description": (
                        f"Slide {idx + 1} {field['label']} image — "
                        f"recommended {rec.get('width', '?')}×{rec.get('height', '?')}px @ 300dpi"
                    ),
                }

            elif field["type"] == "table":
                key = f"slide_{idx}_{field['id']}"
                properties[key] = {
                    "type": "array",
                    "description": (
                        f"Slide {idx + 1} {field['label']} — "
                        f"{field['rows']} rows × {field['columns']} cols"
                    ),
                    "items": {
                        "type": "array",
                        "items": {"type": "string"},
                    },
                }

    return {
        "$schema": "http://json-schema.org/draft-07/schema#",
        "title": f"{manifest['template_name']} Content",
        "type": "object",
        "required": required,
        "properties": properties,
    }


def build_copy_prompt(manifest, template_name, brand_id=""):
    lines = [
        f"# {template_name} — Content Generation Guide\n",
        "## Overview\n",
        f"- **Slides:** {manifest['total_slides']}",
        f"- **Dimensions:** {manifest['slide_dimensions']['width_in']}\" × {manifest['slide_dimensions']['height_in']}\"",
    ]
    if brand_id:
        lines.append(f"- **Brand:** {brand_id}")
    lines += [
        "",
        "Use this guide with voice.yaml and messaging YAML files for brand-aligned content.\n",
        "---\n",
        "## Fields by Slide\n",
    ]

    for slide in manifest["slides"]:
        idx = slide["slide_index"]
        lines.append(f"### Slide {idx + 1} — {slide['layout_name']}\n")

        text_fields = [f for f in slide["fields"] if f["type"] == "text"]
        image_fields = [f for f in slide["fields"] if f["type"] == "image"]
        table_fields = [f for f in slide["fields"] if f["type"] == "table"]

        if text_fields:
            lines.append("| Field | Schema Key | Max Chars | Font | Required | Example |")
            lines.append("|---|---|---|---|---|---|")
            for f in text_fields:
                key = f"slide_{idx}_{f['id']}"
                req = "✓" if f.get("required") else "○"
                ex = (f["example"][:60] + "…") if len(f.get("example", "")) > 60 else f.get("example", "")
                lines.append(
                    f"| {f['label']} | `{key}` | {f['max_chars']} | {f['font_size_pt']}pt | {req} | {ex} |"
                )
            lines.append("")

        for f in image_fields:
            key = f"slide_{idx}_{f['id']}"
            rec = f.get("recommended_px", {})
            lines.append(
                f"**Image: {f['label']}** (`{key}`) — "
                f"{rec.get('width', '?')}×{rec.get('height', '?')}px @ 300dpi\n"
            )

        for f in table_fields:
            key = f"slide_{idx}_{f['id']}"
            lines.append(
                f"**Table: {f['label']}** (`{key}`) — "
                f"{f['rows']} rows × {f['columns']} cols"
                + (" (has header row)" if f.get("has_header_row") else "")
                + "\n"
            )

        lines.append("---\n")

    lines += [
        "## Generation Notes\n",
        "- Match brand voice from `voice.yaml` — characteristics and tone_guidance",
        "- Pull proof points from `messaging/proof-points.yaml` where applicable",
        "- Messaging pillars live in `messaging/messaging-pillars.yaml`",
        "- Respect `max_chars` limits — text that overflows will be cut off in the template",
        "- For image fields: provide a file path or URL; the renderer will resize to recommended_px",
        "- Required fields (✓) must be populated; optional fields (○) can be omitted or left as empty string",
    ]

    return "\n".join(lines)


def extract(pptx_path, output_dir, template_name=None, brand_id="", copy_source=True):
    from pptx import Presentation

    pptx_path = Path(pptx_path).resolve()
    if not pptx_path.exists():
        print(f"Error: file not found: {pptx_path}")
        sys.exit(1)

    if template_name is None:
        template_name = pptx_path.stem

    output_dir = Path(output_dir) / template_name
    output_dir.mkdir(parents=True, exist_ok=True)

    print(f"Extracting template from: {pptx_path.name}")
    print(f"Output directory: {output_dir}")

    prs = Presentation(str(pptx_path))
    slide_width_in = emu_to_inches(prs.slide_width)
    slide_height_in = emu_to_inches(prs.slide_height)

    slides_data = []
    for i, slide in enumerate(prs.slides):
        slide_data = process_slide(slide, i, slide_width_in, slide_height_in)
        slides_data.append(slide_data)
        print(f"  Slide {i + 1}/{len(prs.slides)}: {slide_data['layout_name']} — {len(slide_data['fields'])} fields")

    manifest = build_manifest(str(pptx_path), template_name, slides_data, slide_width_in, slide_height_in)
    content_schema = build_content_schema(manifest)
    copy_prompt = build_copy_prompt(manifest, template_name, brand_id)

    manifest_path = output_dir / "manifest.json"
    schema_path = output_dir / "content-schema.json"
    prompt_path = output_dir / "copy-prompt.md"

    manifest_path.write_text(json.dumps(manifest, indent=2))
    schema_path.write_text(json.dumps(content_schema, indent=2))
    prompt_path.write_text(copy_prompt)

    if copy_source:
        dest_pptx = output_dir / "template.pptx"
        shutil.copy2(str(pptx_path), str(dest_pptx))

    total_fields = sum(len(s["fields"]) for s in slides_data)
    text_fields = sum(
        len([f for f in s["fields"] if f["type"] == "text"]) for s in slides_data
    )

    print(f"\nDone. {len(slides_data)} slides, {total_fields} total fields ({text_fields} text).")
    print(f"  {manifest_path}")
    print(f"  {schema_path}")
    print(f"  {prompt_path}")
    if copy_source:
        print(f"  {output_dir / 'template.pptx'}")

    return str(output_dir)


def main():
    parser = argparse.ArgumentParser(
        description="Extract template specs from a PPTX file for MaC Company Manager"
    )
    parser.add_argument("pptx_path", help="Path to the .pptx file")
    parser.add_argument(
        "--output-dir",
        default=".",
        help="Output directory (a subdirectory named after the template will be created)",
    )
    parser.add_argument(
        "--template-name",
        default=None,
        help="Template name (defaults to the PPTX filename without extension)",
    )
    parser.add_argument(
        "--brand-id",
        default="",
        help="Brand ID to reference in copy-prompt.md",
    )
    parser.add_argument(
        "--no-copy-source",
        action="store_true",
        help="Do not copy the source PPTX to the output directory",
    )
    args = parser.parse_args()

    check_dependency()
    extract(
        pptx_path=args.pptx_path,
        output_dir=args.output_dir,
        template_name=args.template_name,
        brand_id=args.brand_id,
        copy_source=not args.no_copy_source,
    )


if __name__ == "__main__":
    main()
